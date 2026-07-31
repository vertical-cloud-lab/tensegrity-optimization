"""Stamp the IDETC presenter identity onto every deck and every build source.

Marcus Madsen presents the talk, so his name and title belong on the title
slide — and, because the EMC template keeps its title text on the *slide master*
and its presenter block on the *Title Slide layout* rather than on a slide, on
those two parts as well. Otherwise any slide created from the Title Slide layout
silently shows Sterling's EMC 2026 title and affiliation (that is exactly what
Draft 1 slide 14 was doing).

The build sources are stamped too, so a rebuild of the supplement deck or the
background addendum does not reintroduce the EMC text.

Run after any deck rebuild; it is idempotent.

  python presentation/set_presenter.py
"""

from pathlib import Path

from pptx import Presentation

HERE = Path(__file__).parent

PRESENTER = "Marcus Madsen"
ROLE = "Research Assistant"
VENUE = "ASME IDETC-CIE 2026"
WHEN = "August 2026"

# The master title is four runs: two title lines, a spacer, and a gray tagline.
TITLE_RUNS = [
    "\tLet’s build better tensegrity structures",
    "\tfaster, using real drop-test data",
    "\t",
    "+ closed-loop BO and multi-material printing",
]

TARGETS = [
    HERE / "Slide Decks" / "IDETC Tensegrity Slides Draft 1.pptx",
    HERE / "Slide Decks" / "IDETC Supplement Slides (BO block + gap + video + accel).pptx",
    HERE / "Slide Decks" / "IDETC Background Addendum.pptx",
    HERE / "emc2026-idetc-demo.pptx",
    HERE / "emc-bo-block.pptx",          # build source for the supplement deck
    HERE / "emc2026-bare-template.pptx",  # build source for the addendum
]

STALE_TITLE = "high-performing alloys"
STALE_PRESENTERS = ("Sterling Baird", "[Presenter Name]", PRESENTER)


def runs_of(shape):
    """Every run in a shape, in reading order — paragraph breaks included."""
    return [r for p in shape.text_frame.paragraphs for r in p.runs]


def set_runs(runs, texts):
    """Replace run text in place so font, size, and color survive."""
    for run, text in zip(runs, texts):
        run.text = text


def stamp(path):
    prs = Presentation(str(path))
    hits = []

    for master in prs.slide_masters:
        for shape in master.shapes:
            if shape.has_text_frame and STALE_TITLE in shape.text_frame.text:
                set_runs(runs_of(shape), TITLE_RUNS)
                hits.append("master title")
        for layout in master.slide_layouts:
            if layout.name != "Title Slide":
                continue
            for shape in layout.shapes:
                if not shape.has_text_frame:
                    continue
                if not any(n in shape.text_frame.text for n in STALE_PRESENTERS):
                    continue
                runs = runs_of(shape)
                set_runs(runs[:2], [PRESENTER, ROLE])
                if len(runs) >= 4:
                    set_runs(runs[2:4], [VENUE, WHEN])
                hits.append("Title Slide layout")

    # Draft 1 puts the presenter block on the slide itself, as a plain text box.
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and PRESENTER in shape.text_frame.text:
                runs = runs_of(shape)
                if len(runs) >= 2:
                    set_runs(runs[:2], [PRESENTER, ROLE])
                    hits.append("slide presenter block")

    prs.save(str(path))
    print(f"{path.name}: {', '.join(hits) if hits else 'nothing to stamp'}")


if __name__ == "__main__":
    for target in TARGETS:
        if target.exists():
            stamp(target)
