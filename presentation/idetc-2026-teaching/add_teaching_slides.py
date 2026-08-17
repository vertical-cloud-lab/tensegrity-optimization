"""Insert two plain-language teaching slides into the IDETC 2026 deck (issue #94).

Usage: python add_teaching_slides.py <input.pptx> <output.pptx>

- Inserts "Slide A" (how to read one drop) and "Slide B" (smoothing + the
  score) directly after the sensors slide (slide 11), i.e. as slides 12-13.
- Replaces the "Need more information about the drop tests..." placeholder
  paragraph in slide 10's speaker notes with an actual plain-language talk
  track (the rest of that note is kept).
- Adds full speaker-note talk tracks to both new slides.

Everything else in the deck (media, animations, other slides) is untouched.
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent

FIG1 = HERE / "fig1_one_drop_two_parts.png"
FIG2 = HERE / "fig2_smoothing_and_score.png"

TITLE_A = (
    "Each drop tells a two-part story: the jolt going in, "
    "and the ringing that follows."
)
NOTES_A = """How to read the squiggly lines from one real drop.

Left half - the landing itself. The blue curve is the bottom sensor: the jolt the plate delivers to the structure's feet - hundreds of times the pull of gravity, over in about half of one thousandth of a second. The orange curve is the top sensor: the same jolt after it has traveled up through the structure.

Right half - what happens next. The structure keeps swaying, exactly like a bell after it has been struck - here about 560 times per second - and the sway steadily fades as the structure turns that motion into heat.

So each drop gives us two honest readouts: how much of the jolt reaches the top, and how quickly the ringing fades. The fade speed is our most direct sign of energy being soaked up. The ring rate is also a free health check on each print: a print with looser internal tension rings at a measurably different rate."""

TITLE_B = (
    "Every recording gets the same standard treatment, "
    "and each drop boils down to one score."
)
NOTES_B = """Why smooth at all? The raw recording (gray) is dominated by extremely fast wiggles - the metal plate and the sensor crystal shuddering thousands of times per second. Those spikes read over 5,000 G, but they are not the structure moving. Smoothing removes them and leaves the slower push that actually shoves the structure (blue). We do not invent our own smoothing: we use the exact recipe crash-test labs use (an SAE standard called J211), applied identically to every channel of every drop.

An important honest point: the peak you quote depends on how much smoothing you apply - that is exactly why the recipe must be fixed, standard, and identical for every specimen. (It is also why the peaks here read lower than the axis numbers on the previous slide, which was smoothed more lightly for display.)

The score - biggest jolt at the top divided by biggest jolt at the bottom - is a fair way to RANK structures tested the same day on the same pad. Below 1 means the structure softened the jolt; above 1 means the top actually shook harder, which happens when the landing is very sharp. It is NOT "energy absorbed" in joules - no pair of these sensors gives that directly. For the official campaign we therefore also log what a real energy number needs (the falling weight and drop height), record extra quiet time before each jolt, and pair this score with the ring-fade measure from the previous slide. Several of these fixes came out of independent reviews of our earlier settings - the review trail lives in the project repository (issues #86 and #94)."""

SLIDE10_NOTE_REPLACEMENT = """How the test works, in plain terms: the printed structure rides on a flat plate that we raise 60 inches up a pair of rails and release. The plate lands on a stack of felt pads, which turns the landing into a single sharp jolt - like a phone landing on a carpeted floor. Two small motion sensors ride along: one on the plate at the structure's feet, one in a printed pocket at the top point. The recorder wakes itself the instant the jolt begins and takes 1.25 million readings per second for the next 20 thousandths of a second. A drop takes about a minute end to end, so one session yields dozens of repeats."""

PLACEHOLDER_SENTENCE = "Need more information about the drop tests"


def add_teaching_slide(prs, layout, title_text, fig_path, notes_text):
    slide = prs.slides.add_slide(layout)
    # title placeholder, positioned like the rest of the deck's title-only slides
    title = slide.shapes.title
    title.left, title.top = Inches(0.42), Inches(0.4)
    title.width, title.height = Inches(12.5), Inches(0.9)
    title.text = title_text
    # figure: 2480x1120 px -> keep aspect, centered under the title
    w = Inches(12.4)
    h = Inches(12.4 * 1120 / 2480)
    left = (prs.slide_width - w) // 2
    slide.shapes.add_picture(str(fig_path), left, Inches(1.55), width=w, height=h)
    slide.notes_slide.notes_text_frame.text = notes_text
    return slide


def move_slide(prs, from_idx, to_idx):
    sld_ids = prs.slides._sldIdLst
    ids = list(sld_ids)
    el = ids[from_idx]
    sld_ids.remove(el)
    sld_ids.insert(to_idx, el)


def find_slide(prs, title_prefix):
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip().startswith(title_prefix):
                return i
    raise LookupError(title_prefix)


def main(src, dst):
    prs = Presentation(src)
    n = len(prs.slides)
    layout = next(l for l in prs.slide_layouts if l.name == "Title Only")

    # anchor on titles so this works on any revision of the deck
    anchor = find_slide(prs, "We use accelerometers")  # the sensors slide
    add_teaching_slide(prs, layout, TITLE_A, FIG1, NOTES_A)
    add_teaching_slide(prs, layout, TITLE_B, FIG2, NOTES_B)
    # appended at indices n and n+1; move to directly after the sensors slide
    move_slide(prs, n, anchor + 1)
    move_slide(prs, n + 1, anchor + 2)

    # fix the drop-video slide's placeholder note, keep the rest of the note
    s10 = prs.slides[find_slide(prs, "We gather real data")]
    tf = s10.notes_slide.notes_text_frame
    old = tf.text
    if PLACEHOLDER_SENTENCE in old:
        first_para, _, rest = old.partition("\n")
        tf.text = SLIDE10_NOTE_REPLACEMENT + "\n" + rest.lstrip("\n") if rest else SLIDE10_NOTE_REPLACEMENT
    else:
        print("WARNING: slide 10 placeholder note not found; notes left unchanged")

    prs.save(dst)
    print(f"saved {dst}: {len(Presentation(dst).slides)} slides (was {n})")


if __name__ == "__main__":
    main(sys.argv[1], sys.argv[2])
