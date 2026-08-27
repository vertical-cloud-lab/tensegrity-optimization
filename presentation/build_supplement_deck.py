"""Rebuild the IDETC supplement deck: message titles, big type, real video.

Design rules enforced here (the v1 deck broke all of them):
  * one full-sentence message per slide, in the title, <= 2 rendered lines
  * body text >= 24 pt, <= ~160 characters per slide
  * the visual owns the slide; text does not
  * every media slot holds a real clip or figure, not a dashed placeholder

Sterling's EMC 2026 Bayesian-optimization block is carried over untouched
(byte-identical slide XML, animations intact) from the v1 supplement deck.

Usage:  python presentation/build_supplement_deck.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
DECKS = HERE / "Slide Decks"
MEDIA = HERE / "media"
# Read from the pristine EMC block, never from the output deck: SRC == OUT
# made the build non-idempotent, and a second run silently sliced its own
# freshly-written slides 4-14 out in place of the BO explainer.
SRC = HERE / "emc-bo-block.pptx"
OUT = DECKS / "IDETC Supplement Slides (BO block + gap + video + accel).pptx"

# EMC slides to keep from SRC (1-indexed): the whole BO explainer block.
KEEP = list(range(1, 12))

NAVY = RGBColor(0x0E, 0x28, 0x41)
BLUE = RGBColor(0x15, 0x60, 0x82)
ORANGE = RGBColor(0xE9, 0x71, 0x32)
GRAY = RGBColor(0x59, 0x59, 0x59)

SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)
TITLE_TOP = Inches(0.35)
TITLE_H = Inches(1.30)
BODY_TOP = Inches(1.90)
BODY_BOTTOM = Inches(6.95)


def title_only(prs):
    return prs.slides.add_slide(prs.slide_layouts[1])


def set_message(slide, text, size=30):
    """Put the slide's one message in the title placeholder, Doumont-style."""
    ph = slide.shapes.title
    ph.left, ph.top = MARGIN, TITLE_TOP
    ph.width, ph.height = SW - 2 * MARGIN, TITLE_H
    tf = ph.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.color.rgb = NAVY
    return ph


def textbox(slide, left, top, width, height, lines, size=24, color=NAVY,
            bold_first=False, space=10, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(space)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold_first and i == 0
    return box


def credit(slide, text, left=None, top=None, width=None):
    left = MARGIN if left is None else left
    top = Inches(7.02) if top is None else top
    width = SW - 2 * MARGIN if width is None else width
    box = slide.shapes.add_textbox(left, top, width, Inches(0.34))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    for r in tf.paragraphs[0].runs:
        r.font.size = Pt(12)
        r.font.color.rgb = GRAY
        r.font.italic = True
    return box


def fit(box_l, box_t, box_w, box_h, native_w, native_h):
    """Letterbox a native-aspect asset into a box, centered."""
    scale = min(box_w / native_w, box_h / native_h)
    w, h = int(native_w * scale), int(native_h * scale)
    return (int(box_l + (box_w - w) / 2), int(box_t + (box_h - h) / 2), w, h)


def add_video(slide, stem, native, box, poster=None):
    left, top, w, h = fit(*box, *native)
    return slide.shapes.add_movie(
        str(MEDIA / stem), Emu(left), Emu(top), Emu(w), Emu(h),
        poster_frame_image=str(MEDIA / poster) if poster else None,
        mime_type="video/mp4",
    )


def add_image(slide, name, box):
    from PIL import Image

    with Image.open(MEDIA / name) as im:
        native = im.size
    left, top, w, h = fit(*box, *native)
    return slide.shapes.add_picture(str(MEDIA / name), Emu(left), Emu(top),
                                    Emu(w), Emu(h))


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# --------------------------------------------------------------------------
# generated figures
# --------------------------------------------------------------------------

# Callout geometry for the tensegrity anatomy figure. The labels are drawn at
# canvas resolution rather than as PowerPoint text boxes so the leader lines
# can land on specific members; the sizes below put the label text at ~26 pt
# once the figure is placed on the slide.
ANATOMY_CANVAS = (1600, 800)
ANATOMY_PHOTO = (30, 45, 700)  # left, top, height in canvas px
ANATOMY_LABEL_X = 800

# (anchor x, anchor y) in photo-native px; label top; heading; sub-line
ANATOMY_CALLOUTS = [
    ((300, 165), 150, "Struts", "carry compression only"),
    ((360, 263), 320, "No strut touches another",
     "load passes through the cables"),
    ((600, 400), 500, "Cables", "carry tension only"),
]


def build_anatomy_figure():
    """Draw fig-tensegrity-anatomy.png from the 2D teaching-model still."""
    from PIL import Image, ImageDraw, ImageFont

    out = MEDIA / "fig-tensegrity-anatomy.png"
    photo = Image.open(MEDIA / "photo-tensegrity-2d-model.jpg").convert("RGB")
    left, top, height = ANATOMY_PHOTO
    scale = height / photo.height
    photo = photo.resize((int(photo.width * scale), height), Image.LANCZOS)

    canvas = Image.new("RGB", ANATOMY_CANVAS, "white")
    canvas.paste(photo, (left, top))
    draw = ImageDraw.Draw(canvas)

    font_dir = "/usr/share/fonts/truetype/dejavu"
    bold = ImageFont.truetype(f"{font_dir}/DejaVuSans-Bold.ttf", 52)
    plain = ImageFont.truetype(f"{font_dir}/DejaVuSans.ttf", 44)
    orange, navy, gray = (0xE9, 0x71, 0x32), (0x0E, 0x28, 0x41), (0x59,) * 3

    for (ax, ay), ly, head, tail in ANATOMY_CALLOUTS:
        cx, cy = left + ax * scale, top + ay * scale
        draw.line([(ANATOMY_LABEL_X - 30, ly + 34), (cx, cy)], fill=orange,
                  width=7)
        draw.ellipse([cx - 13, cy - 13, cx + 13, cy + 13], fill=orange)
        draw.text((ANATOMY_LABEL_X, ly), head, font=bold, fill=navy)
        draw.text((ANATOMY_LABEL_X, ly + 66), tail, font=plain, fill=gray)
        for text, font in ((head, bold), (tail, plain)):
            if ANATOMY_LABEL_X + draw.textlength(text, font=font) > \
                    ANATOMY_CANVAS[0] - 20:
                raise ValueError(f"callout overruns the canvas: {text!r}")

    canvas.save(out)
    return out


def strip_to_emc_block(prs):
    """Delete every v1 slide except the imported EMC block, preserving XML."""
    id_lst = prs.slides._sldIdLst
    entries = list(id_lst)
    keep_ids = {entries[i - 1] for i in KEEP}
    for entry in entries:
        if entry not in keep_ids:
            rid = entry.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/"
                "relationships}id"
            )
            id_lst.remove(entry)  # unlink first so the rel refcount drops to 0
            prs.part.drop_rel(rid)
    # python-pptx names a new slide part slide<count+1>.xml rather than filling
    # gaps, so the survivors must be renumbered 1..n or new slides collide with
    # them and the package ends up with duplicate zip entries.
    from pptx.opc.packuri import PackURI

    for i, slide in enumerate(prs.slides, start=1):
        slide.part.partname = PackURI(f"/ppt/slides/slide{i}.xml")
    return len(keep_ids)


def move_block_after(prs, block_len, position):
    """Move the leading `block_len` slides to sit after `position` slides."""
    id_lst = prs.slides._sldIdLst
    entries = list(id_lst)
    block, rest = entries[:block_len], entries[block_len:]
    order = rest[:position] + block + rest[position:]
    for e in entries:
        id_lst.remove(e)
    for e in order:
        id_lst.append(e)


# --------------------------------------------------------------------------
# slide builders
# --------------------------------------------------------------------------

def s_readme(prs):
    s = title_only(prs)
    set_message(s, "How to use this supplement (delete before the talk)", 28)
    textbox(s, MARGIN, BODY_TOP, SW - 2 * MARGIN, Inches(4.4), [
        "Slides 8–18 are Sterling's EMC 2026 originals — animations intact. "
        "Copy with “Keep Source Formatting.”",
        "Every video on these slides is a real embedded clip; click to play, "
        "no internet needed.",
        "Working slides (this one, and the last two) are planning aids, not "
        "talk slides.",
    ], size=22, color=NAVY)
    notes(s, "Not a talk slide. Built by presentation/build_supplement_deck.py.")
    return s


def s_hook(prs):
    s = title_only(prs)
    set_message(s, "A tensegrity lander can hit the ground at highway speed, "
                   "bounce, and still get up.")
    add_video(s, "clip-titan-descent.mp4", (1280, 720),
              (Inches(2.15), BODY_TOP, Inches(9.0), Inches(5.0)),
              poster="poster-titan.jpg")
    credit(s, "NASA Super Ball Bot / NIAC Titan mission concept — 20 s excerpt. "
              "Credit: NASA.")
    notes(s, "Open here. Play the clip silently, say nothing for the first "
             "five seconds. Then the scope line: a printed PLA-TPU T3 prism is "
             "our proxy for developing the workflow, not flight hardware.")
    return s


def s_what_is_tensegrity(prs):
    s = title_only(prs)
    set_message(s, "A tensegrity holds itself up with rigid struts that never "
                   "touch, floating in a net of cables.")
    add_video(s, "clip-tensegrity-2d-teaching.mp4", (1280, 718),
              (Inches(2.15), BODY_TOP, Inches(9.0), Inches(5.0)),
              poster="poster-mould-teaching.jpg")
    credit(s, "Steve Mould, “Tensegrity Explained” (youtube.com/watch?v="
              "0onncd0_0-o) — 35 s excerpt, played with sound, used with "
              "on-screen credit.")
    notes(s,
          "BACKGROUND, first beat. Play the whole 35 s with the sound up and "
          "stay quiet — the narration does the teaching, and the 2D model is "
          "the fastest way to make the mechanism obvious to anyone who has "
          "never seen a tensegrity.\n\n"
          "Watch for: the model standing with nothing glued or hinged; the "
          "push; the spring back. That elastic return is the whole reason a "
          "tensegrity is interesting as an energy absorber.\n\n"
          "The earlier 18 s silent crop cut the push-and-recovery off the "
          "end, which threw away the part that teaches. Source: the Box "
          "folder shared on PR #84.")
    return s


def s_tensegrity_anatomy(prs):
    """Freeze the teaching model and name its parts."""
    s = title_only(prs)
    set_message(s, "Every load path is either pure compression in a strut or "
                   "pure tension in a cable — nothing in between.")
    add_image(s, "fig-tensegrity-anatomy.png",
              (MARGIN, BODY_TOP, SW - 2 * MARGIN, Inches(4.85)))
    credit(s, "Still from the same clip — Steve Mould, “Tensegrity "
              "Explained” (youtube.com/watch?v=0onncd0_0-o).")
    notes(s,
          "BACKGROUND, second beat. Freeze the model and name the parts, so "
          "the vocabulary is in place before the tensegrity-inspired caveat "
          "later on.\n\n"
          "Say: this is why a tensegrity can be light and still stiff — no "
          "member ever sees a bending moment, and the cables set the "
          "stiffness. It is also why the geometry matters so much: change a "
          "strut length or a cable pre-tension and the whole response "
          "changes, which is what makes it an optimization problem.\n\n"
          "This is the mechanism visual the mock audience said the grad "
          "student needed in order to tell a tensegrity from a lattice.")
    return s


def s_toy_to_lander(prs):
    s = title_only(prs)
    set_message(s, "NASA's idea started with a baby toy: throw it at the floor, "
                   "nothing breaks — that is a landing robot.")
    add_video(s, "clip-nasa-toy-lander.mp4", (1280, 720),
              (Inches(2.15), BODY_TOP, Inches(9.0), Inches(5.0)),
              poster="poster-nasa-toy.jpg")
    credit(s, "Adrian Agogino, NASA Ames — “NASA 360 Talks: Super Ball Bot” "
              "(youtube.com/watch?v=0eC4A2PXM-U), first 16 s. Credit: NASA.")
    notes(s, "PLAY THIS ONE WITH SOUND — the audio is the slide. He holds a "
             "tensegrity baby toy, says they are made as baby toys because they "
             "are almost impossible to break, throws it at the floor, and lands "
             "on “hey, that's a landing robot.” Say nothing over it; pick up "
             "with our own version of that idea. He reaches “planetary landers” "
             "verbatim a few seconds later in the source video if you ever want "
             "the longer cut (24.5-31 s).")
    return s


def s_reusable(prs):
    s = title_only(prs)
    set_message(s, "Crushable honeycomb absorbs one impact; a tensegrity "
                   "structure survives drop after drop.")
    add_video(s, "clip-superball-3m-drop.mp4", (1280, 722),
              (Inches(2.15), BODY_TOP, Inches(9.0), Inches(5.0)),
              poster="poster-superball.jpg")
    credit(s, "SUPERball v2 3.4 m drop — IEEE Spectrum / NASA Ames "
              "(youtube.com/watch?v=hkzeE6BVNIk), 10 s excerpt.")
    notes(s, "This is the reusability argument Sterling asked for: the same "
             "hardware takes the next drop. It is also why our own campaign "
             "can re-test a specimen instead of consuming it.")
    return s


def s_prior_work(prs):
    s = title_only(prs)
    set_message(s, "Tensegrity impact structures already exist — every one of "
                   "them was designed and assembled by hand.")
    textbox(s, MARGIN, BODY_TOP, Inches(11.5), Inches(4.4), [
        "NASA SUPERball v2 — hand-built, one design, no optimization loop",
        "Davami 2025 — printed tensegrity, characterized but not optimized",
        "Gu & Dotov — tensegrity crutch tip: the application pull is real",
    ], size=26)
    credit(s, "TODO before freeze: verify citations against "
              "manuscript/references.bib.")
    notes(s, "Three names, ten seconds. The point is not a literature review; "
             "it is that the design step is still manual everywhere.")
    return s


def s_gap(prs):
    s = title_only(prs)
    set_message(s, "Nobody closes the loop: print the structure, drop it, and "
                   "let the measured data choose the next design.")
    textbox(s, MARGIN, BODY_TOP, Inches(5.7), Inches(4.2), [
        "TODAY", "Hand-tuned geometry", "Simulation, rarely validated",
        "One-off demonstrations",
    ], size=26, color=GRAY, bold_first=True, space=14)
    textbox(s, Inches(7.0), BODY_TOP, Inches(5.7), Inches(4.2), [
        "THIS WORK", "The drop test is the objective",
        "Trade-offs learned from shock data", "One printer, one tower, one loop",
    ], size=26, color=BLUE, bold_first=True, space=14)
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.05),
                               Inches(3.45), Inches(0.9), Inches(0.6))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ORANGE
    arrow.line.fill.background()
    notes(s, "This is the motivation-to-gap transition. Do not read the "
             "columns; point left, point right, move on.")
    return s


def s_caveat(prs):
    s = title_only(prs)
    set_message(s, "Honest caveat, in the same breath: as printed, these are "
                   "tensegrity-inspired, not true tensegrity.")
    textbox(s, MARGIN, Inches(2.35), Inches(7.4), Inches(3.8), [
        "The black TPU members are not pre-tensioned — and they stretch.",
        "Next: pre-tensioned prints, plus hand-built true-tensegrity "
        "validation specimens.",
    ], size=26, space=26)
    add_image(s, "photo-specimen.jpg",
              (Inches(8.35), Inches(1.95), Inches(4.4), Inches(4.9)))
    credit(s, "Printed PLA–TPU T3 prism, prc1kn — frame from our own 959 fps "
              "drop footage.")
    notes(s, "Say this the moment multi-material printing appears — do not let "
             "the audience discover it in Q&A. Pair it with the path forward "
             "(issue #87 pre-tensioning) so it reads as scope, not weakness.")
    return s


def s_print_timelapse(prs):
    s = title_only(prs)
    set_message(s, "One machine lays down the rigid struts and the stretchy "
                   "cables in a single build.")
    add_video(s, "clip-print-timelapse.mp4", (1120, 720),
              (Inches(6.35), Inches(1.95), Inches(6.45), Inches(4.9)),
              poster="poster-print-timelapse.jpg")
    textbox(s, MARGIN, Inches(2.55), Inches(5.5), Inches(3.8), [
        "Bambu Lab H2D, two materials: PLA struts, TPU cables",
        "No jig, no knots, no post-assembly step",
    ], size=26, space=26)
    credit(s, "Our own print timelapse, TT3_01 — 16 s excerpt "
              "(youtube.com/watch?v=nQNmi-NiL5I, BYU Vertical Cloud Lab).")
    notes(s, "Spoken line: twenty-odd minutes of printing and no assembly step "
             "is what makes a design-per-day loop possible at all. Play it "
             "once, silently.\n\n"
             "Pulled from our own YouTube channel through the lab Raspberry "
             "Pi — the CI runner is bot-blocked by YouTube. Embedded here, so "
             "it needs no internet at the podium.")
    return s


def s_support_removal(prs):
    s = title_only(prs)
    set_message(s, "Printing it in one build is not the same as getting it out "
                   "clean — supports still come off by hand.")
    add_video(s, "clip-support-removal.mp4", (960, 540),
              (Inches(4.1), Inches(2.0), Inches(8.7), Inches(4.8)),
              poster="poster-support-removal.jpg")
    textbox(s, MARGIN, Inches(2.6), Inches(3.3), Inches(3.4), [
        "The honest bottleneck in the loop today",
    ], size=24, color=GRAY)
    credit(s, "Our own footage — manual support removal on a T3 prism "
              "(PR #35, 2026-06-09).")
    notes(s, "Challenge slide, ten seconds. Say the mitigation out loud: "
             "painted manual supports and a 0.4 mm tip cut removal time, and "
             "automating this step is exactly what stands between us and a "
             "self-driving lab. Do not over-dwell; it is a caveat, not a "
             "confession.")
    return s


def s_drop_room(prs):
    s = title_only(prs)
    set_message(s, "The whole experiment takes a quarter of a second — here it "
                   "is in real time.")
    add_video(s, "clip-drop-afar.mp4", (1280, 720),
              (Inches(2.15), BODY_TOP, Inches(9.0), Inches(5.0)),
              poster="poster-drop-afar.jpg")
    credit(s, "Our own footage — drop tower, BYU Smart Materials lab. "
              "Play with sound.")
    notes(s, "PLAY THIS WITH SOUND ON. Seven seconds, real time, whole room. "
             "The bang is the point: it tells the audience how violent a "
             "millisecond-scale event this is before any plot appears. Then "
             "say the standard condition out loud — 60 in onto the same felt "
             "stack every time, repeated per specimen — and move on.")
    return s


def s_drop_phone(prs):
    s = title_only(prs)
    set_message(s, "A phone in slow motion is enough to see the drop — and to "
                   "hear the impact.")
    add_video(s, "clip-drop-phone-audio.mp4", (540, 960),
              (Inches(8.0), Inches(1.85), Inches(4.6), Inches(5.05)),
              poster="poster-drop-phone.jpg")
    textbox(s, MARGIN, Inches(2.45), Inches(6.9), Inches(3.9), [
        "Free, instant, and good enough to catch setup mistakes",
        "How we caught a specimen lifting off and a sensor coming loose",
    ], size=25, space=24)
    credit(s, "Our own footage — specimen n0jdwk, 13 in drop (PR #67). "
              "Play with sound.")
    notes(s, "Play with sound. Roughly eight times slower than real time "
             "(1.6 s of descent for a fall that takes under 0.26 s), which is "
             "the usual 240 fps phone mode. This is the cheap instrument and "
             "it earned its place: phone slow-mo is how we caught the specimen "
             "lifting off the base plate on the bungee-assisted tower, and how "
             "we caught the accelerometer coming loose. It sets up the next "
             "two slides, which say what the fast camera is and is not for.")
    return s


def s_elastic_recovery(prs):
    s = title_only(prs)
    set_message(s, "The top vertex leaves the impact at 0.7 times the speed "
                   "it arrived — it springs back.")
    add_video(s, "clip-drop-highspeed.mp4", (528, 720),
              (Inches(8.2), Inches(1.9), Inches(4.4), Inches(4.95)),
              poster="poster-drop-highspeed.jpg")
    textbox(s, MARGIN, Inches(2.45), Inches(7.1), Inches(3.9), [
        "Elastic re-extension, not permanent collapse",
        "Measured on both validation specimens; both intact afterwards",
    ], size=25, space=24)
    credit(s, "Our own 960 fps footage — specimen 7xadt6, 60 in / 5-felt "
              "validation campaign (PR #86 branch).")
    notes(s, "This is the reusability claim made with our own data instead of "
             "NASA's, and it is squarely in the post-pulse regime the previous "
             "slide said the camera owns. Snap-back at ~0.7x impact speed with "
             "visibly intact specimens is in the 60in-5felts video writeup. "
             "Say it plainly: a crushable absorber would not do this, and it "
             "is why one specimen can serve many drops in the campaign.")
    return s


def s_pipeline(prs):
    s = title_only(prs)
    set_message(s, "One drop becomes one row of data: three numbers and the "
                   "scatter around them.")
    steps = [
        ("Raw", "1.25 MHz\n4 channels"),
        ("Baseline", "pre-trigger\nmedian"),
        ("J211 filter", "CFC-180\nand CFC-1000"),
        ("Metrics", "peak force, SEA,\ncompaction"),
    ]
    x = MARGIN
    w, gap = Inches(2.72), Inches(0.35)
    for i, (head, sub) in enumerate(steps):
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.55),
                                 w, Inches(1.85))
        box.fill.solid()
        box.fill.fore_color.rgb = BLUE if i < 3 else ORANGE
        box.line.fill.background()
        tf = box.text_frame
        tf.word_wrap = True
        tf.text = head
        for r in tf.paragraphs[0].runs:
            r.font.size, r.font.bold = Pt(24), True
        p = tf.add_paragraph()
        p.text = sub
        for r in p.runs:
            r.font.size = Pt(16)
        if i < 3:
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + w + Emu(30000),
                                       Inches(3.25), gap - Emu(60000),
                                       Inches(0.45))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRAY
            arrow.line.fill.background()
        x = x + w + gap
    textbox(s, MARGIN, Inches(4.95), Inches(12.2), Inches(1.8), [
        "Replicate drops give the noise model the optimizer needs — a single "
        "drop is never one point.",
    ], size=24, color=NAVY)
    credit(s, "scripts/analysis/drop_test_*_analysis.py on the PR #86 branch; "
              "raw CSVs committed next to every campaign.")
    notes(s, "Answers issue #94 directly: this is the black box opened. Four "
             "steps, one sentence each. The only judgement calls are step 2 "
             "and step 3 — and both of them bit us, which is the next few "
             "slides.")
    return s


def s_filter_class(prs):
    s = title_only(prs)
    set_message(s, "SAE J211 sets the filter, not us — and the class you pick "
                   "changes the peak by a factor of two.")
    add_image(s, "fig-baseline-and-cfc.png",
              (MARGIN, Inches(1.80), SW - 2 * MARGIN, Inches(4.4)))
    textbox(s, MARGIN, Inches(6.25), Inches(12.2), Inches(0.7), [
        "Same drop, same sensor: raw 567 G · CFC-1000 344 G · CFC-180 245 G",
    ], size=24, color=BLUE)
    credit(s, "Specimen bpx68c, Signal 11 (data/drop-tests/pu-configs/). "
              "Two-pole Butterworth, forward and backward, per SAE J211-1.")
    notes(s, "For issue #94. CFC-180 is a 300 Hz corner, CFC-1000 is 1650 Hz — "
             "the number in the class name times 1.65. Neither is more correct; "
             "CFC-180 is the automotive-crash convention for what a body feels, "
             "CFC-1000 keeps the structural ringing we need to hear the "
             "specimen. That is why our tables report both, and why moving "
             "between them changed a recommendation.")
    return s


def s_two_bands(prs):
    s = title_only(prs)
    set_message(s, "We report both bands because they answer two different "
                   "questions.")
    textbox(s, MARGIN, Inches(2.1), Inches(6.0), Inches(4.3), [
        "CFC-180 · 300 Hz",
        "“What does the payload feel?”",
        "Smooth pulse, stable peak, the number that goes in the force "
        "constraint",
    ], size=25, color=BLUE, bold_first=True, space=18)
    textbox(s, Inches(7.0), Inches(2.1), Inches(5.7), Inches(4.3), [
        "CFC-1000 · 1650 Hz",
        "“What is the structure doing?”",
        "Keeps the 500–550 Hz specimen mode — the part that actually differs "
        "between designs",
    ], size=25, color=ORANGE, bold_first=True, space=18)
    credit(s, "Specimen first mode measured at 519–549 Hz across the ringdown "
              "analyses in this repo.")
    notes(s, "The honest version of the answer to Sterling's question in #94. "
             "We did not abandon CFC-180; we stopped using it as the only "
             "band, because at a 300 Hz corner the specimen's own mode is "
             "filtered away and every design looks alike. If asked which one "
             "feeds the optimizer: the constraint is CFC-180, the "
             "discrimination diagnostics are CFC-1000.")
    return s


def s_adversary(prs):
    s = title_only(prs)
    set_message(s, "We paid an adversary to break our own analysis, and it "
                   "broke it.")
    textbox(s, MARGIN, Inches(1.75), Inches(12.2), Inches(1.5), [
        "Our script assumed the record began at the trigger. It does not — "
        "0.41 ms of it is pre-impact.",
    ], size=24, space=6)
    add_image(s, "fig-baseline-flip.png",
              (Inches(1.55), Inches(2.75), Inches(10.2), Inches(3.35)))
    textbox(s, MARGIN, Inches(6.18), Inches(12.2), Inches(0.75), [
        "Verdict we adopted: none of the four — that sweep could not decide.",
    ], size=24, color=BLUE)
    credit(s, "Edison Scientific adversarial review, task d9092c5a; "
              "recomputation reproduced independently in-repo.")
    notes(s, "This is the credibility slide, and it is the direct answer to "
             "issue #94's “we would like to spot-check you.” Deliver it "
             "without drama: we asked for the analysis to be attacked, four "
             "of our grounds fell, we marked the document superseded rather "
             "than quietly patching it. If a reviewer asks what else it "
             "touched: two sibling analyses share the baseline, both flagged "
             "for re-run.")
    return s


def s_checkable(prs):
    s = title_only(prs)
    set_message(s, "Every number on these slides can be re-derived from "
                   "committed raw data.")
    textbox(s, MARGIN, BODY_TOP, Inches(12.2), Inches(4.4), [
        "Raw CSVs and the analysis script live beside each campaign in the "
        "repo",
        "The filter is a published standard: SAE J211-1 / ISO 6487, channel "
        "class × 1.65 Hz",
        "The adversarial re-analysis, its recomputed tables, and its notebook "
        "are committed too",
    ], size=25, space=18)
    credit(s, "Backup slide: repo paths and the tutorial notebook for issue "
              "#94 are listed in presentation/issue-94-analysis-slides.md.")
    notes(s, "Backup slide for Q&A, and the place to put the Colab link once "
             "the tutorial notebook from issue #94 exists. The spoken version "
             "is one sentence: nothing here is a black box you have to take "
             "on trust.")
    return s


def s_video_capture(prs):
    s = title_only(prs)
    set_message(s, "High-speed video shows the structure spring back and hold "
                   "its shape.")
    add_video(s, "clip-our-slomo-drop.mp4", (720, 720),
              (Inches(7.35), BODY_TOP, Inches(5.4), Inches(5.0)),
              poster="poster-our-drop.jpg")
    textbox(s, MARGIN, Inches(2.35), Inches(6.3), Inches(4.4), [
        "Sony RX100 IV, 959 fps",
        "Rebound, sag and recovery over the next 150 ms",
        "No permanent deformation — the specimen goes back in the tower",
    ], size=26, space=40)
    credit(s, "Our own footage: prc1kn specimen, 60 in drop, 5-felt input "
              "(data/drop-tests/prc1kn-60in-5felt/video/).")
    notes(s, "Play it once, silently, and talk about what they can see: it "
             "bounces, and it comes back intact. That is the reusability "
             "argument in our own data.\n\n"
             "PRESENTER-ONLY, do not volunteer: one frame is 1.04 ms, so the "
             "1.6 ms deceleration itself spans 1-2 frames — the accelerometer "
             "is what measures the pulse. Resolving compression *during* the "
             "pulse would need >=5000 fps DIC "
             "(docs/drop-test-prc1kn-video-analysis.md, section 3). Only bring "
             "this up if someone asks what the camera can resolve.")
    return s


def s_video_processing(prs):
    s = title_only(prs)
    set_message(s, "One clean shock, an elastic rebound, and the brake catches "
                   "the carriage 79 ms later.")
    add_image(s, "fig-video-montage.png",
              (MARGIN, Inches(1.95), SW - 2 * MARGIN, Inches(4.9)))
    credit(s, "prc1kn drop 1, 959.04 fps.")
    notes(s, "Walk it left to right in one sentence: entry, contact, "
             "turnaround, rebound, brake catch, hold. What it proves is rig "
             "physics and specimen integrity — one shock per drop, no plastic "
             "deformation, and the anti-rebound brake catching 150 mm above "
             "the stack.\n\n"
             "PRESENTER-ONLY: contact and turnaround are adjacent frames, so "
             "do not claim the montage resolves the deceleration itself.")
    return s


def s_instrument_split(prs):
    s = title_only(prs)
    set_message(s, "The camera measures the bounce; the accelerometer measures "
                   "the force.")
    add_image(s, "fig-video-impact-zoom.png",
              (MARGIN, Inches(2.00), Inches(7.5), Inches(4.85)))
    textbox(s, Inches(8.35), Inches(2.35), Inches(4.4), Inches(4.3), [
        "Camera → rebound: e* = 0.45, and the structure survives",
        "Accelerometer → peak force, energy absorbed, compaction",
        "Together they cover the whole event",
    ], size=25, space=24)
    credit(s, "prc1kn drop 1, 959 fps; TP4 accelerometer at 125 kHz.")
    notes(s, "One sentence: the video tells us how much of the drop comes "
             "back, the accelerometer tells us what the payload felt. The "
             "coefficient of restitution 0.45 is a free reusability metric — "
             "it needs no pixel scale.\n\n"
             "PRESENTER-ONLY: the objectives come from the accelerometer, not "
             "the video. And do not say the two instruments 'agree on the same "
             "drop' — the prc1kn videos were shot ~5.5 h before that DAQ "
             "campaign, so drop-level pairing is not possible "
             "(docs/drop-test-prc1kn-video-analysis.md).")
    return s


def s_accel_setup(prs):
    s = title_only(prs)
    set_message(s, "A drop tower with an instrumented payload measures the "
                   "force that actually reaches the payload.")
    textbox(s, MARGIN, BODY_TOP, Inches(12.2), Inches(3.6), [
        "Lansmont M23 tower · Dytran 3133A4 accelerometers",
        "TP4 DAQ at 125 kHz — this is the instrument that sees the pulse",
        "Standard test: 60 in drop onto 5 felts, repeated per specimen",
    ], size=26, space=16)
    credit(s, "Setup and standard conditions per the drop-test protocol "
              "(PR #86 branch); 25 kHz is the measured floor for the current "
              "pipeline (issue #89 sample-rate study).")
    notes(s, "Add a photo of the tower here before the talk if one is "
             "available; the text is a stand-in, not a design choice.")
    return s


def s_accel_processing(prs):
    s = title_only(prs)
    set_message(s, "SAE J211 filtering turns raw ringing into the three "
                   "numbers the optimizer actually uses.")
    add_image(s, "fig-impact-zoom-cfc.png",
              (MARGIN, Inches(1.95), SW - 2 * MARGIN, Inches(2.7)))
    textbox(s, MARGIN, Inches(4.95), Inches(12.2), Inches(1.9), [
        "Peak transmitted force  ·  specific energy absorbed  ·  compaction "
        "efficiency",
    ], size=26, color=BLUE)
    credit(s, "Raw vs. CFC-1000 vs. CFC-180 on the same impact window.")
    notes(s, "Do not explain the filter classes. Say: raw data ring at "
             "mounting resonance, the standard filter removes it, and what "
             "survives is the objective function.")
    return s


def s_sensors_lied(prs):
    s = title_only(prs)
    set_message(s, "Our sensors lied to us first — cross-calibration caught a "
                   "5% error before it reached the optimizer.")
    textbox(s, MARGIN, BODY_TOP, Inches(12.2), Inches(3.6), [
        "A mis-entered sensitivity made channel 5 read 0.953× channel 4.",
        "Channel 1 was clipping at the highest drops and we could not see it.",
        "Fix: regress every channel against every other, every campaign.",
    ], size=25, space=16)
    credit(s, "Calibration story documented in PR #74.")
    notes(s, "This is the challenge slide Sterling asked to intersperse. It "
             "buys enormous credibility: we found our own error, and the "
             "cross-check is now routine.")
    return s


def s_challenges_map(prs):
    s = title_only(prs)
    set_message(s, "Working slide: challenges and where each one gets told", 28)
    textbox(s, MARGIN, BODY_TOP, Inches(12.2), Inches(4.6), [
        "Sensor sensitivity error → cross-calibration → own slide (33)",
        "Wrong baseline window → adversarial re-analysis → own slide (34)",
        "Camera and DAQ cover different timescales → each gets its own job "
        "→ slides 24, 27",
        "Supports still removed by hand → painted supports → slide 21",
        "Channel clipping at high drops → headroom check → spoken on slide 28",
        "Print defects confound specimens → replicates → Q&A backup",
        "Cables not in tension → pre-tensioned prints → slide 19",
    ], size=20, space=8)
    notes(s, "Planning aid. Delete before the talk.")
    return s


def s_media_shortlist(prs):
    s = title_only(prs)
    set_message(s, "Working slide: media inventory and what still needs "
                   "shooting", 28)
    textbox(s, MARGIN, BODY_TOP, Inches(12.2), Inches(4.6), [
        "EXTERNAL, EMBEDDED: Titan descent · 2D tensegrity · NASA "
        "toy-to-lander (sound) · SUPERball drop",
        "OURS, EMBEDDED: print timelapse · support removal · whole-room drop "
        "(sound) · phone slow-mo (sound) · 960 fps drop · 7xadt6 snap-back",
        "DUPLICATE: the print timelapse also opens the background addendum — "
        "keep one copy, not both",
        "STILL NEEDED: drop-tower photo, campaign ledger and Pareto-front "
        "figures, Colab link for issue #94",
    ], size=19, space=10)
    notes(s, "Planning aid. Delete before the talk.")
    return s


def main():
    build_anatomy_figure()
    prs = Presentation(str(SRC))
    kept = strip_to_emc_block(prs)
    print(f"kept {kept} EMC slides")

    builders = [
        s_readme, s_hook, s_what_is_tensegrity, s_tensegrity_anatomy,
        s_toy_to_lander, s_reusable,
        s_prior_work, s_gap,
        # <- EMC block gets moved in here (position 8)
        s_caveat, s_print_timelapse, s_support_removal,
        s_drop_room, s_drop_phone, s_video_capture, s_elastic_recovery,
        s_video_processing, s_instrument_split,
        s_accel_setup, s_accel_processing,
        s_pipeline, s_filter_class, s_two_bands, s_sensors_lied, s_adversary,
        s_checkable,
        s_challenges_map, s_media_shortlist,
    ]
    for b in builders:
        b(prs)

    move_block_after(prs, kept, 8)
    prs.save(str(OUT))
    print(f"wrote {OUT} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")


if __name__ == "__main__":
    main()
